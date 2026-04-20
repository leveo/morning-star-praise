import logging
import subprocess
import tempfile
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

from app.config import settings
from app.models import SlideData
from app.services.chinese_service import contains_chinese

logger = logging.getLogger(__name__)

VIDEO_BG_EXTENSIONS = {".mp4", ".webm", ".mov"}
VIDEO_MIME_TYPES = {
    ".mp4": "video/mp4",
    ".webm": "video/webm",
    ".mov": "video/quicktime",
}
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Full <p:timing> template that makes a video shape autoplay + loop in
# slideshow mode. python-pptx's default add_movie timing only configures the
# cMediaNode (click-to-play, no loop). This template adds the <p:seq> +
# <p:cmd type="call" cmd="playFrom(0.0)"> action that PowerPoint's
# "Playback → Start: Automatically + Loop until Stopped" checkboxes generate.
_AUTOPLAY_LOOP_TIMING = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="indefinite"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="afterEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="indefinite" fill="hold"/>
                                      <p:tgtEl>
                                        <p:spTgt spid="{shape_id}"/>
                                      </p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
          <p:video>
            <p:cMediaNode loop="1" mute="1">
              <p:cTn id="7" repeatCount="indefinite" fill="hold" display="0">
                <p:stCondLst>
                  <p:cond delay="indefinite"/>
                </p:stCondLst>
              </p:cTn>
              <p:tgtEl>
                <p:spTgt spid="{shape_id}"/>
              </p:tgtEl>
            </p:cMediaNode>
          </p:video>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


def _is_video_bg(path: Path | None) -> bool:
    return bool(path and path.suffix.lower() in VIDEO_BG_EXTENSIONS)


def _enable_video_autoplay_loop(slide, shape_id: int) -> None:
    """Replace the slide's ``<p:timing>`` with a full autoplay+loop template.

    python-pptx's default ``add_movie`` timing produces a click-to-play,
    no-loop media node. That gets displayed as a static poster image in
    slideshow mode until the user clicks it. For a background video we want
    the slide to start playback the moment it becomes active and keep looping
    — that requires the ``<p:seq>`` + ``<p:cmd ... playFrom(0.0)>`` flow that
    PowerPoint writes when you tick "Start: Automatically" + "Loop until
    Stopped" in the Playback ribbon.
    """
    timing_xml = _AUTOPLAY_LOOP_TIMING.format(shape_id=shape_id).encode("utf-8")
    new_timing = etree.fromstring(timing_xml)

    sld = slide._element
    existing = sld.find(f"{{{_PML_NS}}}timing")
    if existing is not None:
        sld.remove(existing)
    sld.append(new_timing)


def _extract_video_poster(video_path: Path) -> Path | None:
    """Extract the first frame of a video as a JPEG to use as PPT poster frame.

    Without an explicit poster, python-pptx falls back to a 120x55 media icon
    placeholder that PowerPoint then stretches to the slide size — users see a
    blurry icon instead of the video background.
    """
    fd, tmp_name = tempfile.mkstemp(suffix=".jpg", prefix="ppt_poster_")
    import os
    os.close(fd)
    tmp = Path(tmp_name)
    cmd = [
        "ffmpeg", "-y",
        "-ss", "0",
        "-i", str(video_path),
        "-frames:v", "1",
        "-q:v", "2",
        str(tmp),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0 or not tmp.exists() or tmp.stat().st_size == 0:
        logger.warning(
            "Failed to extract poster from %s: %s",
            video_path.name, (result.stderr or "")[-300:],
        )
        tmp.unlink(missing_ok=True)
        return None
    return tmp


def _add_video_background(slide, video_path: Path, slide_width, slide_height):
    """Add a full-slide looping video background to a slide."""
    mime = VIDEO_MIME_TYPES.get(video_path.suffix.lower(), "video/mp4")
    poster = _extract_video_poster(video_path)
    try:
        movie = slide.shapes.add_movie(
            str(video_path),
            0, 0, slide_width, slide_height,
            poster_frame_image=str(poster) if poster else None,
            mime_type=mime,
        )
    finally:
        if poster is not None:
            poster.unlink(missing_ok=True)
    _enable_video_autoplay_loop(slide, movie.shape_id)
    return movie


def _add_background(slide, bg_path: Path | None, slide_width, slide_height) -> None:
    """Add a background to the slide — picture for images, looping movie for videos."""
    if bg_path is None or not bg_path.exists():
        return
    if _is_video_bg(bg_path):
        _add_video_background(slide, bg_path, slide_width, slide_height)
    else:
        slide.shapes.add_picture(str(bg_path), 0, 0, slide_width, slide_height)


def _compute_secondary_size(primary_size: int, secondary_lines: list[str]) -> int:
    """Pick a secondary (translation) font size that keeps long lines on one row.

    Interleaved zh+en is the common case: Chinese primary at 40pt is wide
    enough to fit ~16 chars; an English translation line can easily be 50+
    characters. A flat 0.6 multiplier (= 24pt) is too big and wraps. We scale
    down based on the longest secondary line.
    """
    base = max(int(primary_size * 0.5), 16)
    if not secondary_lines:
        return base
    longest = max(len(ln) for ln in secondary_lines)
    if longest > 55:
        return max(int(primary_size * 0.36), 14)
    if longest > 42:
        return max(int(primary_size * 0.42), 15)
    if longest > 30:
        return max(int(primary_size * 0.46), 15)
    return base


def _set_east_asian_font(run, font_name: str):
    """Set East Asian font via lxml for proper Chinese rendering."""
    rPr = run._r.get_or_add_rPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    ea = rPr.find("a:ea", nsmap)
    if ea is None:
        ea = etree.SubElement(
            rPr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        )
    ea.set("typeface", font_name)


def _padding_palette(padding_style: str) -> dict:
    """Overlay + text colors keyed by padding style. Shared between content
    and title slides so the two can't drift."""
    if padding_style == "light":
        return {
            "overlay": RGBColor(255, 255, 255),
            "primary": RGBColor(0, 0, 0),
            "secondary": RGBColor(60, 60, 60),
            "composer": RGBColor(80, 80, 80),
            # The light backdrop reads heavier than the dark one at the same
            # alpha — bump transparency so the background photo shows through.
            "alpha_text": "55000",
            "alpha_title": "55000",
        }
    return {
        "overlay": RGBColor(0, 0, 0),
        "primary": RGBColor(255, 255, 255),
        "secondary": RGBColor(220, 220, 220),
        "composer": RGBColor(200, 200, 200),
        "alpha_text": "40000",
        "alpha_title": "45000",
    }


def _add_text_with_overlay(
    slide,
    text: str,
    language: str,
    font_size: int | None,
    slide_width,
    slide_height,
    is_title: bool = False,
    *,
    primary_font_size: int | None = None,
    secondary_font_size: int | None = None,
    line_spacing_multiplier: float | None = None,
    padding_style: str = "dark",
):
    """Add semi-transparent overlay and centered text to a slide."""
    palette = _padding_palette(padding_style)
    overlay_rgb = palette["overlay"]
    primary_text_rgb = palette["primary"]
    secondary_text_rgb = palette["secondary"]

    # Uniform padding on all sides
    pad = Inches(0.5)
    overlay_left = pad
    overlay_top = pad
    overlay_width = slide_width - pad * 2
    overlay_height = slide_height - pad * 2

    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        overlay_left,
        overlay_top,
        overlay_width,
        overlay_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = overlay_rgb
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_elem = overlay._element
    for srgb in sp_elem.findall(f".//{{{a_ns}}}srgbClr"):
        if "solidFill" in srgb.getparent().tag:
            alpha = etree.SubElement(srgb, f"{{{a_ns}}}alpha")
            alpha.set("val", palette["alpha_text"])
    overlay.line.fill.background()

    # Text box — fills overlay with small inner margin, vertically centered
    inner_pad = Inches(0.4)
    txBox = slide.shapes.add_textbox(
        overlay_left + inner_pad,
        overlay_top + inner_pad,
        overlay_width - inner_pad * 2,
        overlay_height - inner_pad * 2,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    # Set vertical centering via lxml (bodyPr anchor="ctr")
    body_pr = tf._txBody.find(f"{{{a_ns}}}bodyPr")
    if body_pr is None:
        body_pr = etree.SubElement(tf._txBody, f"{{{a_ns}}}bodyPr")
    body_pr.set("anchor", "ctr")

    # Determine font settings — per-slide override > request-level override > defaults.
    is_zh = language.startswith("zh") or contains_chinese(text)
    if font_size:
        size = font_size
    elif is_title:
        size = settings.TITLE_FONT_SIZE
    elif primary_font_size:
        size = primary_font_size
    elif is_zh:
        size = settings.DEFAULT_FONT_SIZE_ZH
    else:
        size = settings.DEFAULT_FONT_SIZE_EN

    # Detect bilingual layout
    # Stacked: primary block + blank line + secondary block
    # Interleaved: alternating lines of different languages (no blank line)
    sections = text.split("\n\n")
    is_stacked = len(sections) >= 2

    lines = text.split("\n")

    # Detect interleaved: no blank lines, but consecutive lines alternate language
    is_interleaved = False
    if not is_stacked and len(lines) >= 2:
        langs = [contains_chinese(l) for l in lines if l.strip()]
        if len(langs) >= 2:
            alternates = all(langs[i] != langs[i+1] for i in range(len(langs)-1))
            if alternates:
                is_interleaved = True

    # Determine primary language from first non-empty line
    first_line_is_zh = contains_chinese(lines[0]) if lines else is_zh

    # Collect secondary lines so we can size them adaptively. Long English
    # translations under Chinese primaries need smaller type to fit on one row.
    secondary_line_texts: list[str] = []
    if is_stacked:
        in_sec = False
        for ln in lines:
            if ln == "":
                in_sec = True
                continue
            if in_sec and ln.strip():
                secondary_line_texts.append(ln)
    elif is_interleaved:
        for ln in lines:
            if ln.strip() and contains_chinese(ln) != first_line_is_zh:
                secondary_line_texts.append(ln)
    if secondary_font_size:
        secondary_size = secondary_font_size
    else:
        secondary_size = _compute_secondary_size(size, secondary_line_texts)

    zh_spacing = line_spacing_multiplier or 1.5
    en_spacing = line_spacing_multiplier or 1.3

    in_secondary = False
    para_idx = 0

    for line in lines:
        if line == "" and is_stacked:
            in_secondary = True
            p = tf.add_paragraph() if para_idx > 0 else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            spacer_run = p.add_run()
            spacer_run.text = ""
            spacer_run.font.size = Pt(8)
            para_idx += 1
            continue

        p = tf.paragraphs[0] if para_idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.color.rgb = primary_text_rgb

        line_is_zh = contains_chinese(line)

        # Determine if this line is secondary (translation)
        is_secondary_line = False
        if is_stacked:
            is_secondary_line = in_secondary
        elif is_interleaved:
            # In interleaved mode, lines that differ from the first line's language are secondary
            is_secondary_line = (line_is_zh != first_line_is_zh)

        if is_secondary_line:
            run.font.size = Pt(secondary_size)
            run.font.bold = False
            run.font.color.rgb = secondary_text_rgb
            if line_is_zh:
                run.font.name = "PingFang SC"
                _set_east_asian_font(run, "PingFang SC")
            else:
                run.font.name = "Arial"
            p.line_spacing = Pt(secondary_size * (line_spacing_multiplier or 1.3))
        else:
            run.font.size = Pt(size)
            run.font.bold = True
            if line_is_zh:
                run.font.name = "PingFang SC"
                _set_east_asian_font(run, "PingFang SC")
            else:
                run.font.name = "Arial"
            p.line_spacing = Pt(size * (zh_spacing if line_is_zh else en_spacing))

        para_idx += 1


def _add_title_slide(slide, title: str, composer: str, language: str, slide_width, slide_height, *, padding_style: str = "dark"):
    """Add a prominent title slide with song name and composer."""
    is_zh = language.startswith("zh") or contains_chinese(title)
    palette = _padding_palette(padding_style)
    overlay_rgb = palette["overlay"]
    title_rgb = palette["primary"]
    composer_rgb = palette["composer"]

    # Full-slide overlay
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height,
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = overlay_rgb
    sp_elem = overlay._element
    for srgb in sp_elem.findall(f".//{{{a_ns}}}srgbClr"):
        if "solidFill" in srgb.getparent().tag:
            alpha_el = etree.SubElement(srgb, f"{{{a_ns}}}alpha")
            alpha_el.set("val", palette["alpha_title"])
    overlay.line.fill.background()

    # Title text — size based on character count (60-96pt)
    char_count = len(title)
    if char_count <= 4:
        title_size = 96
    elif char_count <= 8:
        title_size = 84
    elif char_count <= 12:
        title_size = 72
    elif char_count <= 20:
        title_size = 66
    else:
        title_size = 60

    title_top = Inches(2.0)
    title_box = slide.shapes.add_textbox(
        Inches(1.0), title_top, slide_width - Inches(2.0), Inches(2.5),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    # Vertical center
    body_pr = tf._txBody.find(f"{{{a_ns}}}bodyPr")
    if body_pr is None:
        body_pr = etree.SubElement(tf._txBody, f"{{{a_ns}}}bodyPr")
    body_pr.set("anchor", "ctr")

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.color.rgb = title_rgb
    run.font.bold = True
    if is_zh:
        run.font.name = "PingFang SC"
        _set_east_asian_font(run, "PingFang SC")
    else:
        run.font.name = "Arial"

    # Composer text — smaller, below title
    if composer:
        comp_box = slide.shapes.add_textbox(
            Inches(1.0), title_top + Inches(2.8), slide_width - Inches(2.0), Inches(1.0),
        )
        tf2 = comp_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = composer
        run2.font.size = Pt(32)
        run2.font.color.rgb = composer_rgb
        if is_zh:
            run2.font.name = "PingFang SC"
            _set_east_asian_font(run2, "PingFang SC")
        else:
            run2.font.name = "Arial"


def _add_page_number(slide, page_num: int, total_pages: int, slide_width):
    """Add page number in top-right corner of a slide."""
    num_box = slide.shapes.add_textbox(
        slide_width - Inches(1.2), Inches(0.2), Inches(1.0), Inches(0.4),
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num}/{total_pages}"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(200, 200, 200)
    run.font.name = "Arial"


def generate_pptx(
    title: str,
    slides: list[SlideData],
    language: str,
    background_paths: list[Path | None],
    composer: str = "",
    show_page_numbers: bool = False,
    primary_font_size: int | None = None,
    secondary_font_size: int | None = None,
    line_spacing_multiplier: float | None = None,
    padding_style: str = "dark",
) -> str:
    """Generate a .pptx file and return the filename."""
    prs = Presentation()
    prs.slide_width = Emu(int(settings.SLIDE_WIDTH_INCHES * 914400))
    prs.slide_height = Emu(int(settings.SLIDE_HEIGHT_INCHES * 914400))

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    blank_layout = prs.slide_layouts[6]  # Blank layout
    total_pages = len(slides)

    # Title slide (no page number)
    slide = prs.slides.add_slide(blank_layout)
    if background_paths:
        _add_background(slide, background_paths[0], slide_width, slide_height)
    _add_title_slide(slide, title, composer, language, slide_width, slide_height, padding_style=padding_style)

    # Content slides
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        bg_idx = (i + 1) % len(background_paths) if background_paths else 0
        bg_path = background_paths[bg_idx] if background_paths else None

        _add_background(slide, bg_path, slide_width, slide_height)

        _add_text_with_overlay(
            slide,
            slide_data.text,
            language,
            slide_data.font_size,
            slide_width,
            slide_height,
            primary_font_size=primary_font_size,
            secondary_font_size=secondary_font_size,
            line_spacing_multiplier=line_spacing_multiplier,
            padding_style=padding_style,
        )

        if show_page_numbers:
            _add_page_number(slide, i + 1, total_pages, slide_width)

    # Save
    filename = f"worship_{uuid.uuid4().hex[:8]}.pptx"
    output_path = settings.OUTPUT_DIR / filename
    prs.save(str(output_path))

    return filename
